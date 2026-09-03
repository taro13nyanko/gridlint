"""Build the 10-slide submission deck.

    python deck/make_deck.py

Writes deck/Gridlint.pptx. Run deck/export_pdf.ps1 afterwards to get a PDF.
The visual language matches the product: teal for what the code proves,
red for the number that is wrong.
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "deck" / "Gridlint.pptx"
SHOTS = ROOT / "deck" / "shots"

INK = RGBColor(0x16, 0x21, 0x1F)
INK2 = RGBColor(0x55, 0x63, 0x5F)
INK3 = RGBColor(0x8B, 0x97, 0x93)
BRAND = RGBColor(0x0F, 0x76, 0x6E)
BRAND_WASH = RGBColor(0xE6, 0xF4, 0xF1)
DANGER = RGBColor(0xD1, 0x34, 0x4B)
DANGER_WASH = RGBColor(0xFD, 0xEC, 0xEF)
PANEL = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF7, 0xF8, 0xF7)
LINE = RGBColor(0xE2, 0xE7, 0xE5)

SANS = "Segoe UI"
MONO = "Consolas"

W, H = Inches(13.333), Inches(7.5)


def new_deck() -> Presentation:
    p = Presentation()
    p.slide_width, p.slide_height = W, H
    return p


def slide(prs, *, bg=BG):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg
    rect.line.fill.background()
    rect.shadow.inherit = False
    return s


def text(s, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.15):
    """runs: list of (text, size, bold, colour) or (text, size, bold, colour, font)."""
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for run in runs:
        body, size, bold, colour = run[:4]
        font_name = run[4] if len(run) > 4 else SANS
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.alignment = align
        para.line_spacing = spacing
        para.space_after = Pt(6)
        r = para.add_run()
        r.text = body
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = colour
        r.font.name = font_name
    return box


def card(s, x, y, w, h, *, fill=PANEL, line=LINE, radius=True):
    shape = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    shape.shadow.inherit = False
    if radius:
        shape.adjustments[0] = 0.06
    shape.text_frame.text = ""
    return shape


def eyebrow(s, label, y=0.55):
    text(s, 0.85, y, 10, 0.3, [(label.upper(), 12, True, BRAND)])


def title(s, t, y=0.95, size=34, colour=INK):
    text(s, 0.85, y, 11.6, 1.1, [(t, size, True, colour)], spacing=1.05)


def footer(s, n):
    text(s, 0.85, 6.95, 8, 0.3, [("Gridlint  ·  AI Builders Hackathon 2026", 10, False, INK3)])
    text(s, 11.5, 6.95, 1, 0.3, [(str(n), 10, False, INK3)], align=PP_ALIGN.RIGHT)


def picture(s, name, x, y, w):
    p = SHOTS / name
    if not p.exists():
        card(s, x, y, w, w * 0.62)
        text(s, x + 0.2, y + 0.2, w - 0.4, 0.4, [(f"[{name} missing]", 12, False, INK3)])
        return None
    pic = s.shapes.add_picture(str(p), Inches(x), Inches(y), width=Inches(w))
    return pic


# --------------------------------------------------------------------- slides

def s1_title(prs):
    s = slide(prs, bg=PANEL)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.28), H)
    band.fill.solid(); band.fill.fore_color.rgb = BRAND; band.line.fill.background()
    band.shadow.inherit = False

    text(s, 1.1, 1.5, 11, 0.4, [("SPREADSHEET REVIEW, MEASURED IN MONEY", 13, True, BRAND)])
    text(s, 1.1, 2.0, 11, 1.0, [("Gridlint", 62, True, INK)], spacing=1.0)
    text(s, 1.1, 3.1, 11.2, 1.4, [
        ("Your board deck says 38.6 months of runway.", 30, False, INK2),
        ("It has 5.2.", 30, True, DANGER),
    ], spacing=1.15)
    text(s, 1.1, 4.9, 11, 0.9, [
        ("Gridlint reads the formulas in a workbook, recalculates every cell, and finds the "
         "mistakes that quietly change the answer — ranked by how much money each one moves.", 16, False, INK2)])
    text(s, 1.1, 6.3, 11, 0.8, [
        ("Yusuke Ota  ·  University of Tokyo  ·  solo entry", 14, True, INK),
        ("github.com/taro13nyanko/gridlint", 13, False, BRAND),
    ])


def s2_problem(prs):
    s = slide(prs)
    eyebrow(s, "The problem")
    title(s, "A broken spreadsheet does not look broken.")
    text(s, 0.85, 2.0, 6.2, 2.2, [
        ("A SUM that stops one row short still returns a number. It still formats as "
         "currency. It still lands in the board deck.", 17, False, INK2),
        ("Field audits find errors in 24% to 94% of real business spreadsheets. Panko found "
         "at least one in 94% of 88 audited workbooks. (EuSpRIG)", 17, False, INK2),
    ])

    card(s, 7.4, 1.95, 5.1, 2.35, fill=DANGER_WASH, line=DANGER_WASH)
    text(s, 7.75, 2.2, 4.4, 1.9, [
        ("=SUM(C11:C14)", 22, True, DANGER, MONO),
        ("Row 15 is “Contractors”, $185,000 a month.", 14, False, INK2),
        ("The total is right there in the deck, and it is short every single month.", 14, False, INK2),
    ])

    text(s, 0.85, 4.6, 11.6, 0.4, [("What people use today", 15, True, INK)])
    cols = [
        ("PerfectXL · Operis OAK · Spreadsheet Detective",
         "Desktop add-ins, $249–$2,000 per year. They hand an analyst 400 warnings."),
        ("Excel's own error checking",
         "Green triangles for a handful of patterns. No ranking, no impact, no fix."),
        ("A second pair of eyes",
         "The most common control in practice, and the least reliable one."),
    ]
    for i, (h, b) in enumerate(cols):
        x = 0.85 + i * 3.95
        card(s, x, 5.05, 3.7, 1.6)
        text(s, x + 0.25, 5.28, 3.2, 1.2, [(h, 13, True, INK), (b, 12, False, INK2)])

    text(s, 0.85, 6.5, 11.6, 0.4,
         [("None of them tells you which warning is worth reading, or what the number becomes once it is fixed.",
           14, True, INK)])
    footer(s, 2)


def s3_solution(prs):
    s = slide(prs)
    eyebrow(s, "The solution")
    title(s, "Rank by money. Then prove the fix.")
    steps = [
        ("1", "Read the formulas",
         "Every formula is parsed into a syntax tree and linked into a dependency graph."),
        ("2", "Recalculate the workbook",
         "Gridlint's own engine recomputes every cell — then checks itself against the values Excel saved."),
        ("3", "Price each defect",
         "The fix is applied to a copy, the workbook is recomputed, and the difference is measured."),
        ("4", "Say it in plain English",
         "A model writes the note. It never decides what is wrong and never produces a number."),
    ]
    for i, (n, h, b) in enumerate(steps):
        x = 0.85 + i * 3.0
        card(s, x, 2.15, 2.75, 2.5)
        badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.25), Inches(2.4),
                                   Inches(0.42), Inches(0.42))
        badge.fill.solid(); badge.fill.fore_color.rgb = BRAND_WASH
        badge.line.fill.background(); badge.shadow.inherit = False
        tf = badge.text_frame; tf.text = n
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].runs[0].font.size = Pt(14)
        tf.paragraphs[0].runs[0].font.bold = True
        tf.paragraphs[0].runs[0].font.color.rgb = BRAND
        text(s, x + 0.25, 3.0, 2.3, 1.5, [(h, 15, True, INK), (b, 12, False, INK2)])

    card(s, 0.85, 5.0, 11.65, 1.6, fill=PANEL)
    text(s, 1.15, 5.2, 11.0, 1.3, [
        ("The output is not a list of warnings. It is a sentence with a number in it:", 14, False, INK2),
        ("“One SUM leaves out the Contractors row. Runway goes from 38.6 months to 5.2, "
         "and gross margin from +6.3% to −30%.”", 17, True, INK),
    ])
    footer(s, 3)


def s4_users(prs):
    s = slide(prs)
    eyebrow(s, "Who it is for")
    title(s, "Whoever signs off on the number.")
    people = [
        ("Finance and ops teams", "Month-end close, board packs, budget re-forecasts. The ritual is monthly and the file is the same one every time."),
        ("Analysts and consultants", "They inherit a model built by somebody who left. Gridlint is the first hour of that handover."),
        ("Grant and programme admins", "Public money, spreadsheet reporting, an auditor at the end of it."),
    ]
    for i, (h, b) in enumerate(people):
        x = 0.85 + i * 3.95
        card(s, x, 2.05, 3.7, 1.85)
        text(s, x + 0.25, 2.3, 3.2, 1.5, [(h, 16, True, INK), (b, 12.5, False, INK2)])

    text(s, 0.85, 4.25, 11.6, 0.4, [("Pricing", 15, True, INK)])
    tiers = [("Free", "$0", "One workbook · every rule"),
             ("Team", "$29 / editor / month", "Unlimited workbooks · history · shareable reports · corrected file"),
             ("Close", "$99 / month", "Scheduled re-checks · CI check for committed workbooks")]
    for i, (name, price, what) in enumerate(tiers):
        x = 0.85 + i * 3.95
        card(s, x, 4.65, 3.7, 1.55, line=BRAND if i == 1 else LINE)
        text(s, x + 0.25, 4.88, 3.2, 1.2,
             [(name, 13, True, INK2), (price, 20, True, INK), (what, 11.5, False, INK2)])
    text(s, 0.85, 6.4, 11.6, 0.4,
         [("It is a subscription because the mistake comes back every month, with the next version of the file.",
           14, True, INK)])
    footer(s, 4)


def s5_product(prs):
    s = slide(prs)
    eyebrow(s, "The product")
    title(s, "One screen: what is wrong, and what it costs.")
    picture(s, "report.png", 0.85, 1.95, 8.4)
    card(s, 9.55, 1.95, 2.95, 4.4)
    text(s, 9.8, 2.2, 2.5, 4.0, [
        ("Ranked by measured impact", 13, True, INK),
        ("The biggest number is the one that moves the most money, not the first one found.", 11.5, False, INK2),
        ("Grouped", 13, True, INK),
        ("Twelve monthly totals with the same mistake are one finding, not twelve.", 11.5, False, INK2),
        ("Verified fix", 13, True, INK),
        ("“30 cells change, no new errors” — because it was actually run.", 11.5, False, INK2),
        ("The grid, flagged", 13, True, INK),
        ("Toggle formulas and the defect is visible in place.", 11.5, False, INK2),
    ])
    footer(s, 5)


def s6_architecture(prs):
    s = slide(prs)
    eyebrow(s, "Technical architecture")
    title(s, "An Excel engine, written from scratch.")
    boxes = [
        (".xlsx", "loaded twice:\nformulas, and the\nvalues Excel cached"),
        ("tokenizer\n→ parser", "single pass, precedence\nclimbing, full AST"),
        ("dependency\ngraph", "topological order,\nTarjan cycle detection"),
        ("recalculate", "42 functions,\nExcel's coercion rules"),
        ("13 rules", "structural, correctness,\nhygiene — no model"),
        ("price + rank", "apply fix → recompute\n→ diff → money"),
    ]
    x = 0.85
    for i, (h, b) in enumerate(boxes):
        w = 1.83
        card(s, x, 2.1, w, 1.5, fill=PANEL, line=BRAND if i in (2, 3) else LINE)
        text(s, x + 0.15, 2.28, w - 0.3, 1.2,
             [(h, 13, True, BRAND if i in (2, 3) else INK), (b, 10.5, False, INK2)])
        if i < len(boxes) - 1:
            text(s, x + w + 0.01, 2.65, 0.3, 0.3, [("→", 14, True, INK3)], align=PP_ALIGN.CENTER)
        x += w + 0.13

    card(s, 0.85, 3.9, 5.7, 2.6, fill=BRAND_WASH, line=BRAND_WASH)
    text(s, 1.15, 4.15, 5.1, 2.2, [
        ("The self-check", 16, True, BRAND),
        ("Every workbook already stores the values its app last computed. Gridlint recomputes "
         "them and compares.", 13, False, INK2),
        ("That single comparison does two jobs: it validates the engine, and it detects a file "
         "saved with calculation switched off — the numbers on screen no longer matching the "
         "formulas underneath.", 13, False, INK2),
        ("Below 99.5% agreement, Gridlint says so and withholds the money figures rather than "
         "guessing.", 13, True, INK),
    ])
    card(s, 6.8, 3.9, 5.7, 2.6)
    text(s, 7.1, 4.15, 5.1, 2.2, [
        ("Engineering decisions", 16, True, INK),
        ("Evaluation never recurses into other cells: the graph is sorted, then evaluated in "
         "order. A 3,000-cell chain cannot blow the stack, and there is a test for it.", 12.5, False, INK2),
        ("Ranges are clamped to the used area, so =SUM(A:A) does not expand to a million "
         "nodes.", 12.5, False, INK2),
        ("Impacts are never summed — they overlap. An odd mistake for this product to make.",
         12.5, False, INK2),
    ])
    footer(s, 6)


def s7_ai(prs):
    s = slide(prs)
    eyebrow(s, "Where the AI is, and where it is not")
    title(s, "The model writes. The code decides.")
    rows = [
        ("Parse formulas, build the graph", "Code"),
        ("Recalculate the workbook", "Code"),
        ("Decide whether something is a defect", "Code"),
        ("Measure what a fix changes", "Code"),
        ("Explain the defect to a non-expert", "Model, fenced"),
        ("Draft a repair where no mechanical fix exists", "Model, executed"),
    ]
    y = 2.05
    for i, (job, who) in enumerate(rows):
        is_model = who.startswith("Model")
        card(s, 0.85, y, 5.9, 0.52, fill=DANGER_WASH if is_model else PANEL,
             line=DANGER_WASH if is_model else LINE)
        text(s, 1.1, y + 0.13, 3.6, 0.3, [(job, 12.5, False, INK)])
        text(s, 4.75, y + 0.13, 1.85, 0.3,
             [(who, 12, True, DANGER if is_model else BRAND)], align=PP_ALIGN.RIGHT)
        y += 0.6

    card(s, 7.05, 2.05, 5.45, 3.35, fill=PANEL, line=BRAND)
    text(s, 7.35, 2.3, 4.9, 3.0, [
        ("Two fences, both tested", 16, True, BRAND),
        ("1 · Number guard", 13, True, INK),
        ("Every figure in a written note must already appear in the evidence the detector "
         "measured. Anything else is a number the model invented, and the whole sentence is "
         "thrown away.", 12, False, INK2),
        ("2 · Executed repairs", 13, True, INK),
        ("A proposed formula is parsed, run through the engine, and diffed. If it introduces "
         "one new error cell anywhere, it is rejected and never shown.", 12, False, INK2),
    ])
    card(s, 0.85, 5.65, 11.65, 0.95, fill=BRAND_WASH, line=BRAND_WASH)
    text(s, 1.15, 5.87, 11.0, 0.7, [
        ("Delete the model entirely and Gridlint still finds every defect and still measures "
         "every impact. The model makes the report readable; it does not make it true.", 15, True, INK)])
    footer(s, 7)


def s8_proof(prs):
    s = slide(prs)
    eyebrow(s, "Proof")
    title(s, "Three numbers, all reproducible from the repository.")
    stats = [
        ("93 / 93", "formulas where the engine matches Excel exactly",
         "A conformance workbook generated by Excel itself: −2^2 = 4, 2^3^2 = 64, "
         "ROUND half away from zero, text ignored in a range but coerced as an argument."),
        ("141 / 141", "planted defects found by the right rule",
         "Six kinds of defect injected into twelve clean workbooks, every mutant recalculated "
         "by Excel, then scored on whether the right rule named the right cell."),
        ("0", "findings on 12 clean workbooks, 1,368 formulas",
         "Blank spacer rows, subtotals rolled into a grand total, cross-sheet assumptions, "
         "VLOOKUP tables. A checker that cries wolf is one nobody keeps installed."),
    ]
    for i, (n, label, body) in enumerate(stats):
        x = 0.85 + i * 3.95
        card(s, x, 2.05, 3.7, 2.6)
        text(s, x + 0.25, 2.3, 3.2, 2.2, [
            (n, 34, True, BRAND), (label, 13, True, INK), (body, 11.5, False, INK2)])

    card(s, 0.85, 4.95, 11.65, 1.5, fill=DANGER_WASH, line=DANGER_WASH)
    text(s, 1.15, 5.18, 11.0, 1.2, [
        ("The benchmark earned its keep during the build.", 15, True, DANGER),
        ("It caught a real bug in Gridlint's own arithmetic: ROUND used a relative epsilon to "
         "correct binary representation error, which rounded 3,806,241.4967 up to 3,806,242. "
         "It now rounds in decimal. That is what a benchmark is for.", 13, False, INK2),
    ])
    footer(s, 8)


def s9_impact(prs):
    s = slide(prs)
    eyebrow(s, "Impact")
    title(s, "The check moves from “someone should look at this”\nto something that runs.", size=30)
    items = [
        ("Before a number leaves the building",
         "Drop the file in, read one sentence, decide. The whole check takes under a second on a "
         "123-formula model."),
        ("Every month, on the same file",
         "A workspace keeps the workbook and its history, so the second month is a re-check, not "
         "a re-review."),
        ("In CI, on a workbook in a repository",
         "gridlint check model.xlsx --fail-on critical exits non-zero. A pull request that breaks "
         "a total fails the build, the way it would for code."),
        ("As a shareable report",
         "A read-only link is how a reviewer says “here is what I found” to whoever owns the file, "
         "without sending the file back."),
    ]
    y = 2.35
    for h, b in items:
        card(s, 0.85, y, 11.65, 0.98)
        text(s, 1.15, y + 0.16, 4.3, 0.7, [(h, 14, True, INK)])
        text(s, 5.6, y + 0.16, 6.6, 0.7, [(b, 12.5, False, INK2)])
        y += 1.1
    footer(s, 9)


def s10_roadmap(prs):
    s = slide(prs)
    eyebrow(s, "Roadmap")
    title(s, "What is built, and what is next.")
    built = [
        "Formula engine, 42 functions, 93/93 against Excel",
        "Dependency graph, recalculation, cycle detection",
        "13 rules, impact measured by recalculation",
        "Web app: workspaces, history, shared reports, corrected download",
        "CLI, CI workflow, Docker, 113 tests",
    ]
    nxt = [
        "Google Sheets: read a file directly instead of exporting",
        "More of Excel: array formulas, XLOOKUP, SUMIFS, date functions",
        "Rule packs per domain (grant reporting, payroll, VAT)",
        "Watch mode: re-check on every save, diff against last month",
        "A public benchmark corpus so other tools can be measured against it",
    ]
    card(s, 0.85, 2.05, 5.7, 3.6, fill=BRAND_WASH, line=BRAND_WASH)
    text(s, 1.15, 2.3, 5.1, 3.2,
         [("Working today", 17, True, BRAND)] + [("·  " + b, 13.5, False, INK2) for b in built])
    card(s, 6.8, 2.05, 5.7, 3.6)
    text(s, 7.1, 2.3, 5.1, 3.2,
         [("Next", 17, True, INK)] + [("·  " + b, 13.5, False, INK2) for b in nxt])

    text(s, 0.85, 5.95, 11.65, 0.9, [
        ("I pay for university by fixing other people's spreadsheets. Every single one had a "
         "silent formula bug in it.", 16, True, INK),
        ("github.com/taro13nyanko/gridlint  ·  MIT", 13, False, BRAND),
    ])
    footer(s, 10)


def main() -> int:
    prs = new_deck()
    for fn in (s1_title, s2_problem, s3_solution, s4_users, s5_product,
               s6_architecture, s7_ai, s8_proof, s9_impact, s10_roadmap):
        fn(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"wrote {OUT} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
